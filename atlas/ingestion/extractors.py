"""Text extractors for the filesystem ingestion source (ADR-0033).

Each extractor turns a file into plain text. New formats are additive: add a
function and register its extensions in ``EXTRACTORS``. Heavy parsers (pypdf,
BeautifulSoup) are imported lazily so a missing optional dependency only affects
its own file type.

``extract(path)`` returns the text, or ``None`` when the file has no usable text
(e.g. a scanned/image-only PDF), so the caller can skip it. Scanned PDFs are a
future OCRService concern.
"""

from __future__ import annotations

import csv
import json
import re
from pathlib import Path
from typing import Any

CONTENT_TYPES = {
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".pdf": "application/pdf",
    ".html": "text/html",
    ".htm": "text/html",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".csv": "text/csv",
    ".json": "application/json",
}


def content_type_for(path: Path) -> str:
    return CONTENT_TYPES.get(path.suffix.lower(), "text/plain")


def _extract_text(path: Path) -> str | None:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    return text or None


def _extract_pdf(path: Path) -> str | None:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(p.strip() for p in pages if p.strip()).strip()
    return text or None  # None => no text layer (scanned) -> skip for future OCR


def _extract_docx(path: Path) -> str | None:
    """Word document → paragraph text + flattened table cells (python-docx)."""
    from docx import Document

    doc = Document(str(path))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    text = "\n".join(parts).strip()
    return text or None


def _extract_pptx(path: Path) -> str | None:
    """Presentation → text of every shape on every slide (python-pptx)."""
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        slide_parts.append(line)
        if slide_parts:
            parts.append(f"# Slide {i}\n" + "\n".join(slide_parts))
    text = "\n\n".join(parts).strip()
    return text or None


def _extract_xlsx(path: Path) -> str | None:
    """Spreadsheet → per-sheet, tab-separated rows (openpyxl, values only)."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for ws in wb.worksheets:
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = ["" if v is None else str(v) for v in row]
                if any(c.strip() for c in cells):
                    rows.append("\t".join(cells).rstrip())
            if rows:
                parts.append(f"# Sheet: {ws.title}\n" + "\n".join(rows))
    finally:
        wb.close()
    text = "\n\n".join(parts).strip()
    return text or None


def _extract_csv(path: Path) -> str | None:
    """CSV → tab-separated rows (dialect-sniffed, robust to odd delimiters)."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    if not raw.strip():
        return None
    try:
        dialect = csv.Sniffer().sniff(raw[:4096])
    except csv.Error:
        dialect = csv.excel
    rows = ["\t".join(row) for row in csv.reader(raw.splitlines(), dialect)]
    text = "\n".join(r for r in rows if r.strip()).strip()
    return text or None


def _extract_json(path: Path) -> str | None:
    """JSON → pretty-printed text (falls back to raw if it doesn't parse)."""
    raw = path.read_text(encoding="utf-8", errors="replace")
    try:
        data = json.loads(raw)
    except (ValueError, TypeError):
        return raw.strip() or None
    return json.dumps(data, indent=2, ensure_ascii=False, sort_keys=False).strip() or None


def html_to_text(html: str) -> str | None:
    """Strip an HTML document down to readable plain text (script/style removed)."""
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator="\n")
    # Collapse blank lines and surrounding whitespace.
    lines = [line.strip() for line in text.splitlines()]
    cleaned = "\n".join(line for line in lines if line).strip()
    return cleaned or None


# Structural boilerplate tags that never carry article body text.
_BOILERPLATE_TAGS = (
    "script", "style", "noscript", "nav", "header", "footer", "aside",
    "form", "button", "iframe", "svg", "figure", "figcaption",
)
# id/class/role substrings that mark non-article regions across publishers.
_BOILERPLATE_HINTS = (
    "nav", "menu", "sidebar", "footer", "header", "banner", "cookie",
    "consent", "subscribe", "advert", "promo", "breadcrumb", "related",
    "recommend", "social", "share", "masthead", "toolbar", "skip", "search",
    "newsletter", "paywall", "modal", "popup", "notice",
)
# Below this the main-content heuristic is considered too thin to trust; the
# caller falls back to the naive whole-page strip (small/simple pages still work).
_MIN_MAIN_CHARS = 200


def _trafilatura_main_text(html: str) -> str | None:
    """Best-effort article extraction via trafilatura (optional dependency)."""
    try:
        import trafilatura
    except Exception:  # noqa: BLE001 - optional; degrade to the heuristic
        return None
    try:
        text = trafilatura.extract(
            html,
            include_comments=False,
            include_tables=True,
            favor_recall=True,
            no_fallback=False,
        )
    except Exception:  # noqa: BLE001 - never let extraction crash the reader
        return None
    text = (text or "").strip()
    return text or None


def _heuristic_main_text(html: str) -> str | None:
    """Deterministic main-content extraction: drop chrome, keep the densest body.

    Removes boilerplate tags and any element whose id/class/role marks it as
    navigation/promo/cookie chrome, then selects the container with the most
    paragraph text (preferring semantic ``<article>``/``<main>``/``[role=main]``).
    Returns ``None`` when nothing text-rich is found so the caller can fall back.
    """
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(list(_BOILERPLATE_TAGS)):
        tag.decompose()
    for el in list(soup.find_all(True)):
        cls = el.get("class")
        cls_str = " ".join(cls) if isinstance(cls, list) else (cls or "")
        attrs = f"{cls_str} {el.get('id') or ''} {el.get('role') or ''}".lower()
        if attrs.strip() and any(h in attrs for h in _BOILERPLATE_HINTS):
            el.decompose()

    def _para_score(node: Any) -> int:
        return sum(len(p.get_text(strip=True)) for p in node.find_all("p"))

    best = None
    best_score = 0
    # Prefer semantic containers, then fall back to the densest div/section.
    semantic = soup.find_all(["article", "main"]) + soup.find_all(attrs={"role": "main"})
    for node in semantic:
        score = _para_score(node)
        if score > best_score:
            best, best_score = node, score
    if best is None or best_score < _MIN_MAIN_CHARS:
        for node in soup.find_all(["div", "section"]):
            score = _para_score(node)
            if score > best_score:
                best, best_score = node, score
    if best is None or best_score < _MIN_MAIN_CHARS:
        return None
    text = best.get_text(separator="\n")
    lines = [line.strip() for line in text.splitlines()]
    cleaned = "\n".join(line for line in lines if line).strip()
    return cleaned or None


def html_to_main_text(html: str) -> str | None:
    """Extract the main article text from HTML, dropping nav/boilerplate.

    Publisher landing pages (IEEE/Springer/ScienceDirect) otherwise flatten into a
    boilerplate blob with no recognizable article body → zero claims. This tries
    trafilatura (best publisher coverage), then a deterministic boilerplate-stripping
    heuristic, and finally the naive :func:`html_to_text` so simple pages still work.
    """
    if not html or not html.strip():
        return None
    return (
        _trafilatura_main_text(html)
        or _heuristic_main_text(html)
        or html_to_text(html)
    )


# Markers that a page is a paywall / login gate rather than an article body.
_PAYWALL_RE = re.compile(
    r"\b(?:subscribe to (?:continue|view|read)|sign in to (?:continue|view|read|access)|"
    r"log in to (?:view|read|access|continue)|purchase (?:this )?(?:article|pdf|access)|"
    r"buy (?:this )?article|get access|access through your institution|"
    r"institutional (?:login|access|sign)|check (?:if you have )?access|"
    r"already (?:a subscriber|have an account)|create an account to)\b",
    re.IGNORECASE,
)


def looks_paywalled(text: str) -> bool:
    """Heuristic: does this normalized text read like a paywall/login gate?"""
    return bool(text) and _PAYWALL_RE.search(text) is not None


def _extract_html(path: Path) -> str | None:
    return html_to_main_text(path.read_text(encoding="utf-8", errors="replace"))


EXTRACTORS = {
    ".txt": _extract_text,
    ".md": _extract_text,
    ".pdf": _extract_pdf,
    ".html": _extract_html,
    ".htm": _extract_html,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".csv": _extract_csv,
    ".json": _extract_json,
}


def supported_extensions() -> list[str]:
    """Sorted list of file extensions with a registered extractor."""
    return sorted(EXTRACTORS)


def extract(path: Path) -> str | None:
    """Extract plain text from ``path`` by file type, or None if unsupported/empty."""
    extractor = EXTRACTORS.get(path.suffix.lower())
    if extractor is None:
        return None
    return extractor(path)
