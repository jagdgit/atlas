"""Tests for the expanded document extractors (S13 Document Reader).

Hermetic: files are synthesised in a tmp dir with the same libraries used to read
them, then extracted back to text.
"""

from __future__ import annotations

import json

import pytest

from atlas.ingestion.extractors import (
    content_type_for,
    extract,
    supported_extensions,
)


def test_supported_extensions_covers_s13_set():
    exts = set(supported_extensions())
    assert {".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".json", ".md", ".txt",
            ".html"} <= exts


def test_extract_txt_and_md(tmp_path):
    p = tmp_path / "note.md"
    p.write_text("# Title\n\nHello Atlas.", encoding="utf-8")
    assert "Hello Atlas." in extract(p)
    assert content_type_for(p) == "text/markdown"


def test_extract_csv(tmp_path):
    p = tmp_path / "data.csv"
    p.write_text("name,score\nAtlas,42\nJagadesh,99\n", encoding="utf-8")
    text = extract(p)
    assert "Atlas" in text and "42" in text
    assert content_type_for(p) == "text/csv"


def test_extract_json_pretty(tmp_path):
    p = tmp_path / "cfg.json"
    p.write_text(json.dumps({"b": 1, "a": [1, 2]}), encoding="utf-8")
    text = extract(p)
    assert '"b": 1' in text  # pretty-printed
    assert content_type_for(p) == "application/json"


def test_extract_json_malformed_falls_back_to_raw(tmp_path):
    p = tmp_path / "bad.json"
    p.write_text("{not valid json", encoding="utf-8")
    assert "not valid json" in extract(p)


def test_extract_docx(tmp_path):
    from docx import Document

    p = tmp_path / "doc.docx"
    d = Document()
    d.add_paragraph("First paragraph.")
    table = d.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Cell A"
    table.rows[0].cells[1].text = "Cell B"
    d.save(str(p))

    text = extract(p)
    assert "First paragraph." in text
    assert "Cell A | Cell B" in text


def test_extract_pptx(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    p = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Slide headline"
    prs.save(str(p))

    text = extract(p)
    assert "Slide headline" in text
    assert "Slide 1" in text


def test_extract_xlsx(tmp_path):
    from openpyxl import Workbook

    p = tmp_path / "sheet.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["name", "value"])
    ws.append(["soiling", 3.7])
    wb.save(str(p))

    text = extract(p)
    assert "Sheet: Data" in text
    assert "soiling" in text and "3.7" in text


def test_extract_unsupported_returns_none(tmp_path):
    p = tmp_path / "thing.xyz"
    p.write_text("data", encoding="utf-8")
    assert extract(p) is None
